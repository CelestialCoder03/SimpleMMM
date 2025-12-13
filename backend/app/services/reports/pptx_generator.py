"""PowerPoint report generator for MMM results."""

import io
import logging
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


class PPTXReportGenerator:
    """
    PowerPoint report generator for Marketing Mix Model results.

    Creates professional presentations with model insights,
    metrics, contributions, and recommendations.
    """

    # Brand colors
    PRIMARY_COLOR = RgbColor(102, 126, 234)  # Blue-purple
    SECONDARY_COLOR = RgbColor(118, 75, 162)  # Purple
    ACCENT_COLOR = RgbColor(76, 175, 80)  # Green
    TEXT_COLOR = RgbColor(51, 51, 51)  # Dark gray
    LIGHT_BG = RgbColor(248, 249, 250)  # Light gray

    def __init__(
        self,
        model_name: str,
        project_name: str,
        model_result: dict[str, Any],
        contributions: dict[str, list[float]] | None = None,
    ):
        """
        Initialize PowerPoint generator.

        Args:
            model_name: Name of the model.
            project_name: Name of the project.
            model_result: Model result data.
            contributions: Channel contribution data.
        """
        self.model_name = model_name
        self.project_name = project_name
        self.model_result = model_result
        self.contributions = contributions or {}

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def generate(self) -> bytes:
        """Generate the PowerPoint presentation."""
        self._add_title_slide()
        self._add_executive_summary_slide()
        self._add_model_metrics_slide()
        self._add_contributions_slide()
        self._add_coefficients_slide()
        self._add_recommendations_slide()
        self._add_appendix_slide()

        # Save to bytes
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.getvalue()

    def _add_title_slide(self):
        """Add title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            self.prs.slide_width,
            self.prs.slide_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.PRIMARY_COLOR
        shape.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Marketing Mix Model Results"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{self.project_name} | {self.model_name}"
        p.font.size = Pt(24)
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        # Date
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%B %d, %Y")
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(200, 200, 200)
        p.alignment = PP_ALIGN.CENTER

    def _add_executive_summary_slide(self):
        """Add executive summary slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Executive Summary")

        metrics = self.model_result.get("metrics", {})
        coefficients = self.model_result.get("coefficients", {})

        # Key metrics section
        y_pos = Inches(1.5)

        # R-squared
        r2 = metrics.get("r_squared", 0)
        self._add_metric_box(slide, Inches(0.5), y_pos, "Model Fit (R²)", f"{r2:.1%}")

        # MAPE
        mape = metrics.get("mape", 0)
        self._add_metric_box(slide, Inches(3.5), y_pos, "MAPE", f"{mape:.1%}")

        # Channels
        n_channels = len([k for k in coefficients if k not in ["intercept", "base"]])
        self._add_metric_box(slide, Inches(6.5), y_pos, "Channels", str(n_channels))

        # Observations
        n_obs = metrics.get("n_observations", 0)
        self._add_metric_box(slide, Inches(9.5), y_pos, "Observations", str(n_obs))

        # Key findings text
        findings_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(3.5))
        tf = findings_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = "Key Findings"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_COLOR

        # Add findings
        findings = self._generate_findings()
        for finding in findings[:5]:
            p = tf.add_paragraph()
            p.text = f"• {finding}"
            p.font.size = Pt(14)
            p.font.color.rgb = self.TEXT_COLOR
            p.space_before = Pt(8)

    def _add_model_metrics_slide(self):
        """Add detailed model metrics slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Model Performance Metrics")

        metrics = self.model_result.get("metrics", {})

        # Create metrics table
        rows = [
            ("R-Squared", f"{metrics.get('r_squared', 0):.4f}"),
            ("Adjusted R-Squared", f"{metrics.get('adj_r_squared', 0):.4f}"),
            ("RMSE", f"{metrics.get('rmse', 0):,.2f}"),
            ("MAE", f"{metrics.get('mae', 0):,.2f}"),
            ("MAPE", f"{metrics.get('mape', 0):.2%}"),
            ("AIC", f"{metrics.get('aic', 0):,.0f}"),
            ("BIC", f"{metrics.get('bic', 0):,.0f}"),
        ]

        table = slide.shapes.add_table(
            len(rows) + 1,
            2,
            Inches(2),
            Inches(1.8),
            Inches(9),
            Inches(0.4 * (len(rows) + 1)),
        ).table

        # Header
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"

        for i, (metric, value) in enumerate(rows):
            table.cell(i + 1, 0).text = metric
            table.cell(i + 1, 1).text = value

        # Style table
        for row_idx in range(len(rows) + 1):
            for col_idx in range(2):
                cell = table.cell(row_idx, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                if row_idx == 0:
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.PRIMARY_COLOR
                    cell.text_frame.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)

    def _add_contributions_slide(self):
        """Add channel contributions slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Channel Contributions")

        if not self.contributions:
            self._add_no_data_message(slide)
            return

        # Calculate total contributions
        contrib_totals = {}
        for channel, values in self.contributions.items():
            if channel not in ["intercept", "base", "trend"]:
                contrib_totals[channel] = sum(values) if values else 0

        total = sum(contrib_totals.values())

        # Sort by contribution
        sorted_contribs = sorted(contrib_totals.items(), key=lambda x: x[1], reverse=True)

        # Create table
        rows = sorted_contribs[:10]  # Top 10
        table = slide.shapes.add_table(
            len(rows) + 1,
            3,
            Inches(1),
            Inches(1.8),
            Inches(11),
            Inches(0.4 * (len(rows) + 1)),
        ).table

        # Set column widths
        table.columns[0].width = Inches(5)
        table.columns[1].width = Inches(3)
        table.columns[2].width = Inches(3)

        # Header
        table.cell(0, 0).text = "Channel"
        table.cell(0, 1).text = "Total Contribution"
        table.cell(0, 2).text = "Share %"

        for i, (channel, contrib) in enumerate(rows):
            share = contrib / total * 100 if total > 0 else 0
            table.cell(i + 1, 0).text = channel
            table.cell(i + 1, 1).text = f"{contrib:,.0f}"
            table.cell(i + 1, 2).text = f"{share:.1f}%"

        # Style table
        for row_idx in range(len(rows) + 1):
            for col_idx in range(3):
                cell = table.cell(row_idx, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(11)
                if row_idx == 0:
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.PRIMARY_COLOR
                    cell.text_frame.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)

    def _add_coefficients_slide(self):
        """Add model coefficients slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Model Coefficients")

        coefficients = self.model_result.get("coefficients", {})
        std_errors = self.model_result.get("std_errors", {})

        if not coefficients:
            self._add_no_data_message(slide)
            return

        # Filter and sort coefficients
        coef_items = [(k, v) for k, v in coefficients.items() if k not in ["intercept"]]
        coef_items.sort(key=lambda x: abs(x[1]), reverse=True)

        # Create table
        rows = coef_items[:12]  # Top 12
        table = slide.shapes.add_table(
            len(rows) + 1,
            3,
            Inches(1.5),
            Inches(1.8),
            Inches(10),
            Inches(0.35 * (len(rows) + 1)),
        ).table

        # Header
        table.cell(0, 0).text = "Variable"
        table.cell(0, 1).text = "Coefficient"
        table.cell(0, 2).text = "Std Error"

        for i, (var, coef) in enumerate(rows):
            se = std_errors.get(var, 0)
            table.cell(i + 1, 0).text = var
            table.cell(i + 1, 1).text = f"{coef:.4f}"
            table.cell(i + 1, 2).text = f"{se:.4f}" if se else "-"

        # Style table
        for row_idx in range(len(rows) + 1):
            for col_idx in range(3):
                cell = table.cell(row_idx, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                if row_idx == 0:
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.PRIMARY_COLOR
                    cell.text_frame.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)

    def _add_recommendations_slide(self):
        """Add recommendations slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Recommendations")

        recommendations = self._generate_recommendations()

        y_pos = Inches(1.8)
        for i, rec in enumerate(recommendations[:6]):
            rec_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12.333), Inches(0.8))
            tf = rec_box.text_frame
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.text = f"{i + 1}. {rec}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.TEXT_COLOR

            y_pos += Inches(0.9)

    def _add_appendix_slide(self):
        """Add appendix slide with technical details."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        self._add_slide_title(slide, "Appendix: Technical Details")

        details = [
            f"Model Type: {self.model_result.get('model_type', 'Unknown')}",
            f"Training Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
            f"Observations: {self.model_result.get('metrics', {}).get('n_observations', 'N/A')}",
            f"Features: {len(self.model_result.get('coefficients', {}))}",
            "Generated by Marketing Mix Model Platform",
        ]

        y_pos = Inches(1.8)
        for detail in details:
            box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12.333), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = detail
            p.font.size = Pt(14)
            p.font.color.rgb = self.TEXT_COLOR
            y_pos += Inches(0.5)

    def _add_slide_title(self, slide, title: str):
        """Add title to a slide."""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY_COLOR

    def _add_metric_box(self, slide, x, y, label: str, value: str):
        """Add a metric box to the slide."""
        width = Inches(2.5)
        height = Inches(1.5)

        # Box shape
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.LIGHT_BG
        shape.line.color.rgb = self.PRIMARY_COLOR

        # Value
        value_box = slide.shapes.add_textbox(x, y + Inches(0.2), width, Inches(0.7))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.PRIMARY_COLOR
        p.alignment = PP_ALIGN.CENTER

        # Label
        label_box = slide.shapes.add_textbox(x, y + Inches(0.9), width, Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = self.TEXT_COLOR
        p.alignment = PP_ALIGN.CENTER

    def _add_no_data_message(self, slide):
        """Add no data message to slide."""
        msg_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.333), Inches(1))
        tf = msg_box.text_frame
        p = tf.paragraphs[0]
        p.text = "No data available for this section"
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RgbColor(128, 128, 128)
        p.alignment = PP_ALIGN.CENTER

    def _generate_findings(self) -> list[str]:
        """Generate key findings from model results."""
        findings = []
        metrics = self.model_result.get("metrics", {})
        coefficients = self.model_result.get("coefficients", {})

        r2 = metrics.get("r_squared", 0)
        if r2 > 0.8:
            findings.append(f"Strong model fit with R² of {r2:.1%}")
        elif r2 > 0.6:
            findings.append(f"Moderate model fit with R² of {r2:.1%}")
        else:
            findings.append(f"Model explains {r2:.1%} of variance in target")

        # Top contributors
        if self.contributions:
            contrib_totals = {k: sum(v) for k, v in self.contributions.items() if k not in ["intercept", "base"]}
            if contrib_totals:
                top_channel = max(contrib_totals, key=contrib_totals.get)
                findings.append(f"{top_channel} is the top contributing channel")

        # Positive/negative coefficients
        positive = [k for k, v in coefficients.items() if v > 0 and k != "intercept"]
        if positive:
            findings.append(f"{len(positive)} channels show positive impact on target")

        findings.append("Model captures key marketing dynamics")

        return findings

    def _generate_recommendations(self) -> list[str]:
        """Generate recommendations from model results."""
        recommendations = []
        coefficients = self.model_result.get("coefficients", {})

        # Sort by coefficient magnitude
        sorted_coefs = sorted(
            [(k, v) for k, v in coefficients.items() if k not in ["intercept", "base"]],
            key=lambda x: x[1],
            reverse=True,
        )

        if sorted_coefs:
            top = sorted_coefs[0]
            recommendations.append(f"Increase investment in {top[0]} - highest positive impact")

        if len(sorted_coefs) > 1:
            second = sorted_coefs[1]
            if second[1] > 0:
                recommendations.append(f"Maintain or grow {second[0]} spend - strong ROI indicator")

        # Low performers
        negatives = [(k, v) for k, v in sorted_coefs if v < 0]
        if negatives:
            worst = negatives[-1]
            recommendations.append(f"Review {worst[0]} strategy - showing diminishing returns")

        recommendations.extend(
            [
                "Conduct regular model refresh with new data",
                "Test incremental budget shifts before major reallocations",
                "Monitor external factors that may impact model accuracy",
            ]
        )

        return recommendations


def generate_pptx_report(
    model_name: str,
    project_name: str,
    model_result: dict[str, Any],
    contributions: dict[str, list[float]] | None = None,
) -> bytes:
    """
    Convenience function to generate PowerPoint report.

    Args:
        model_name: Name of the model.
        project_name: Name of the project.
        model_result: Model result data.
        contributions: Channel contribution data.

    Returns:
        PowerPoint file as bytes.
    """
    generator = PPTXReportGenerator(
        model_name=model_name,
        project_name=project_name,
        model_result=model_result,
        contributions=contributions,
    )
    return generator.generate()
